from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import subprocess
import os
from pdf2image import convert_from_path

def generate_presentation(slide_contents, pptx_path, config=None):
    """
    Generate a PowerPoint presentation from slide contents, which can be:
    1. A list of SlideItem objects
    2. A SlideChunk object
    3. The raw parsed content from the paste.txt
    
    Parameters:
    - slide_contents: Slide content in one of the above formats
    - pptx_path: Path where the presentation will be saved
    - config: Optional configuration parameters
    
    Returns:
    - Path to the saved presentation
    """
    
/Users/indreshgoswami/Downloads/pdfvideo/Pdf_to_Video_Generator/presentation.pptx
    print("Creating enhanced slides...")
    prs = Presentation()
    
    # Handle different possible input formats
    if hasattr(slide_contents, 'slides'):
        # This is a SlideChunk object
        slides = slide_contents.slides
        theme_colors = slide_contents.theme_colors if slide_contents.theme_colors else {}

    elif isinstance(slide_contents, list) and all(hasattr(item, 'title') for item in slide_contents):
        # This is a list of SlideItem objects
        slides = slide_contents
        theme_colors = {}

    elif isinstance(slide_contents, list) and len(slide_contents) > 0 and 'slides' in dir(slide_contents[0]):
        # This might be a list containing one SlideChunk object
        slides = slide_contents[0].slides
        theme_colors = slide_contents[0].theme_colors if slide_contents[0].theme_colors else {}

    else:
        slides = slide_contents
        theme_colors = {}
        
        if hasattr(slide_contents, 'theme_colors'):
            theme_colors = slide_contents.theme_colors
    
    # Ensure we have default theme colors if not provided
    if not theme_colors:
        theme_colors = {
            "primary": "#1F497D",
            "secondary": "#4F81BD", 
            "accent": "#C0504D",
            "background": "#FFFFFF",
            "text": "#000000"
        }
    
    # Strip '#' prefix from colors if present
    for key in theme_colors:
        if isinstance(theme_colors[key], str) and theme_colors[key].startswith('#'):
            theme_colors[key] = theme_colors[key][1:]
    
    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    # Determine presentation title from first slide
    presentation_title = "Presentation"
    if slides and hasattr(slides[0], 'title'):
        first_title = slides[0].title
        if "Introduction" in first_title and "to" in first_title:
            presentation_title = first_title.split("to")[1].strip()
        else:
            presentation_title = first_title
    
    title.text = presentation_title
    subtitle.text = "A Comprehensive Guide"
    
    # Style the title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["secondary"])
    
    # Process each slide
    for i, slide_item in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        # Get slide title
        slide_title = f"Slide {i+1}"
        if hasattr(slide_item, 'title'):
            slide_title = slide_item.title
        
        # Set the title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
        
        # Set the content
        content_frame = content_placeholder.text_frame
        content_frame.clear()
        
        # Add main content
        slide_content = "Content"
        if hasattr(slide_item, 'content'):
            slide_content = slide_item.content
            
        p = content_frame.add_paragraph()
        p.text = slide_content
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor.from_string(theme_colors["text"])
        
        # Add key points as bullets
        key_points = []
        if hasattr(slide_item, 'key_points'):
            key_points = slide_item.key_points
            
        if key_points:
            content_frame.add_paragraph().text = ""  # Add spacing
            for point in key_points:
                bullet_p = content_frame.add_paragraph()
                bullet_p.text = point
                bullet_p.font.size = Pt(20)
                bullet_p.level = 1
                bullet_p.font.color.rgb = RGBColor.from_string(theme_colors["secondary"])
    
   
    
    # Add a final slide
    final_slide = prs.slides.add_slide(prs.slide_layouts[2])
    final_title = final_slide.shapes.title
    final_content = final_slide.placeholders[1]
    
    final_title.text = "Thank You!"
    final_title.text_frame.paragraphs[0].font.size = Pt(40)
    final_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
    
    final_p = final_content.text_frame.add_paragraph()
    final_p.text = "Any questions?"
    final_p.font.size = Pt(32)
    final_p.font.color.rgb = RGBColor.from_string(theme_colors["accent"])
    
    # Save the presentation
    prs.save(pptx_path)
    print(f"✅ Enhanced slides created and saved to {pptx_path}")
    return pptx_path

def slides_to_images(ppt_path, output_folder):
    subprocess.run([
        '/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'pdf', ppt_path, '--outdir', output_folder
    ], check=True)
    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    return [img.save(os.path.join(output_folder, f"slide_{i}.png"), 'PNG') or os.path.join(output_folder, f"slide_{i}.png")
            for i, img in enumerate(convert_from_path(pdf_path, dpi=200))]
