import fitz  # PyMuPDF
from openai import OpenAI
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips
from pdf2image import convert_from_path
from pydantic import BaseModel, ValidationError
import subprocess
import tempfile
import os
import json
import textwrap
from dotenv import load_dotenv

load_dotenv()

client = OpenAI(
    base_url=os.getenv("ENDPOINT"),
    api_key=os.getenv("TOKEN"),
)


# Data models
class SlideItem(BaseModel):
    title: str
    content: str
    key_points: List[str] = Field(default_factory=list)
    image_prompt: Optional[str] = None

class ShortVideoSegment(BaseModel):
    title: str
    content: str
    script: str
    duration: float = 60.0  # Target duration in seconds

class SlideChunk(BaseModel):
    slides: List[SlideItem]
    voice_over_script: str
    short_segments: List[ShortVideoSegment] = Field(default_factory=list)
    theme_colors: Optional[Dict[str, str]] = None

class VideoConfig(BaseModel):
    theme: str = "professional"  # professional, creative, minimal
    presenter_type: str = "human"  # human, cartoon, none
    language: str = "en"
    voice_style: str = "neutral"  # neutral, enthusiastic, formal
    include_background_music: bool = True
    resolution: str = "1080p"
    aspect_ratio: str = "16:9"
    animation_level: str = "moderate"  # none, subtle, moderate, dynamic

# Step 1: Extract PDF Content
def extract_text_from_pdf(pdf_path):
    print("Extracting text from PDF...")
    doc = fitz.open(pdf_path)

    # Extract text and any images
    text = ""
    images = []

    for page_num, page in enumerate(doc):
        text += page.get_text()

        # Extract images if needed
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            # Could save or process images here if needed

    print("✅ Text extracted.")
    return text

# Step 2: Chunk & Summarize
def chunk_text(text, chunk_size=100000):
    print("Chunking text...")
    chunks = textwrap.wrap(text, chunk_size, break_long_words=False)
    print(f"✅ Created {len(chunks)} chunks.")
    return chunks

def generate_chunk_content(chunks, config):
    print("Summarizing chunks with OpenAI...")
    combined_content = "\n\n".join(chunks)

    theme_desc = {
        "professional": "formal, corporate style with clean design",
        "creative": "vibrant, engaging style with dynamic elements",
        "minimal": "clean, simple style with focus on key content"
    }.get(config.theme, "professional style")

    voice_desc = {
        "neutral": "balanced and clear",
        "enthusiastic": "energetic and engaging",
        "formal": "serious and professional"
    }.get(config.voice_style, "clear and professional")

    prompt = (
        f"Generate a structured presentation based on the following content. "
        f"Use a {theme_desc} visual approach and a {voice_desc} tone for narration.\n\n"
        "Create a JSON with these keys:\n"
        "1. 'slides': list of objects with 'title', 'content', 'key_points' (list of bullet points), "
        "and 'image_prompt' (a description for generating a relevant image)\n"
        "2. 'voice_over_script': professional narration script covering all slides\n"
        "3. 'short_segments': 3-5 stand-alone segments for short-form videos (under 2 minutes each) "
        "with 'title', 'content', 'script', and 'duration' fields\n"
        "4. 'theme_colors': suggested color scheme (primary, secondary, accent, background, text)\n\n"
        f"Content:\n{combined_content}\n\n"
        "Respond with valid JSON only. Keep all content factual and based on the input material. "
        "Ensure 'voice_over_script' is a single string, not a list."
    )

    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=16000
    ).choices[0].message.content.strip()

    # Clean up response to get valid JSON
    if response.startswith("```json"):
        response = response.lstrip("```json").rstrip("```").strip()
    elif response.startswith("```"):
        response = response.lstrip("```").rstrip("```").strip()

    try:
        parsed_response = json.loads(response)

        # Fix: If voice_over_script is a list, join it into a single string
        if isinstance(parsed_response.get('voice_over_script'), list):
            parsed_response['voice_over_script'] = "\n\n".join(parsed_response['voice_over_script'])

        validated_chunk = SlideChunk(**parsed_response)
    except (json.JSONDecodeError, ValidationError) as e:
        print(f"Parsing error: {e}\nResponse was: {response}")
        raise

    print("✅ Summarization complete.")
    return validated_chunk

# Step 3: Audio Generation
# def generate_audio(script, filepath, language="en", voice_style="neutral"):
#     print(f"Generating audio in {language}...")

#     # Could be expanded to use different TTS engines based on voice_style
#     tts = gTTS(script, lang=language)
#     tts.save(filepath)

#     # Optionally process audio for better quality (normalize, etc.)
#     # This would require additional audio processing libraries

#     print("✅ Audio generated.")
#     return filepath

# Step 3: Generate Audio from Script
def generate_audio(script, output_file):
    tts = gTTS(script)
    tts.save(output_file)

# Utility function to split audio for short segments
def split_audio_script(full_script, segments):
    print("Generating audio for short segments...")
    segment_audios = []

    for i, segment in enumerate(segments):
        audio_path = f"temp/short_segment_{i}.mp3"
        generate_audio(segment.script, audio_path)
        segment_audios.append(audio_path)

    print(f"✅ Generated {len(segment_audios)} segment audio files.")
    return segment_audios

# Step 4: Enhanced Slides Generation
def generate_presentation(slide_content, pptx_path, config):
    print("Creating enhanced slides...")
    prs = Presentation()

    # Set theme colors based on content analysis
    theme_colors = slide_content.theme_colors or {
        "primary": "1F497D",    # Dark blue
        "secondary": "4F81BD",  # Medium blue
        "accent": "C0504D",     # Red accent
        "background": "FFFFFF", # White background
        "text": "000000"        # Black text
    }

    # Strip '#' prefix from colors if present
    for key in theme_colors:
        if theme_colors[key].startswith('#'):
            theme_colors[key] = theme_colors[key][1:]  # Remove the '#' character

    # Apply slide theme and layout adjustments based on config
    for i, slide_item in enumerate(slide_content.slides):
        # Add title slide for first slide
        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = slide_item.title
            subtitle.text = "Generated Presentation"

            # Apply styling to title slide
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])
        else:
            # Content slides
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
            title = slide.shapes.title
            content_placeholder = slide.placeholders[1]

            title.text = slide_item.title

            # Style the title
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(theme_colors["primary"])

            # Clear existing content and add formatted content
            content_frame = content_placeholder.text_frame
            content_frame.clear()

            # Add main content paragraph
            p = content_frame.add_paragraph()
            p.text = slide_item.content
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor.from_string(theme_colors["text"])

            # Add key points as bullet points if available
            if slide_item.key_points:
                content_frame.add_paragraph().text = ""  # Add spacing

                for point in slide_item.key_points:
                    bullet_p = content_frame.add_paragraph()
                    bullet_p.text = point
                    bullet_p.font.size = Pt(20)
                    bullet_p.level = 1  # Make it a bullet point

            # Could add image placeholder here - in production version,
            # would integrate with image generation API using image_prompt

    prs.save(pptx_path)
    print("✅ Enhanced slides created.")
    return pptx_path


# -------------------- $$$$$$$ -----------------# -------------------- $$$$$$$ -----------------
# # Pydantic schema for validation
# class SlideChunk(BaseModel):
#     slides: list[str]
#     voice_over_script: str

# # Step 1: Extract Text from PDF
# def extract_text_from_pdf(pdf_path):
#     doc = fitz.open(pdf_path)
#     return "".join(page.get_text() for page in doc)

# # Step 2: Chunk and summarize text for slides and voiceover
# def chunk_text(text, max_chunk_chars=2500):
#     paragraphs = textwrap.wrap(text, max_chunk_chars, break_long_words=False, replace_whitespace=False)
#     return paragraphs

# def generate_chunk_content(chunk):
#     prompt = (
#         "Return JSON with two fields: 'slides' (list of concise bullet points summarizing the content) "
#         "and 'voice_over_script' (a detailed voice-over narration). "
#         "Format: {\"slides\": [...], \"voice_over_script\": \"...\"}. "
#         "Only output raw JSON, no markdown formatting like triple backticks.\n\n"
#         f"Content:\n{chunk}\n\nJSON Output:"
#     )
#     response = client.chat.completions.create(
#         model="gpt-4o",
#         messages=[{"role": "user", "content": prompt}],
#         temperature=0.2,
#         max_tokens=4000
#     ).choices[0].message.content.strip()

#     if response.startswith("```json"):
#         response = response.lstrip("```json").rstrip("```").strip()
#     elif response.startswith("```"):
#         response = response.lstrip("```").rstrip("```").strip()

#     try:
#         parsed = json.loads(response)
#         validated = SlideChunk(**parsed)
#     except (json.JSONDecodeError, ValidationError) as e:
#         raise ValueError(f"GPT format error: {e}\nRaw: {response}")

#     return validated

# # Step 3: Generate Audio from Script
# def generate_audio(script, output_file):
#     tts = gTTS(script)
#     tts.save(output_file)

# # Step 4: Generate Multi-Slide Presentation
# def generate_presentation(slide_chunks, ppt_file):
#     prs = Presentation()
#     for i, chunk in enumerate(slide_chunks):
#         slide = prs.slides.add_slide(prs.slide_layouts[1])
#         slide.shapes.title.text = f"Slide {i+1}"
#         frame = slide.placeholders[1].text_frame
#         frame.clear()
#         for point in chunk.slides:
#             p = frame.add_paragraph()
#             p.text = point
#             p.font.size = Pt(24)
#             p.space_after = Pt(12)
#     prs.save(ppt_file)

# Step 5: Convert Slides to Images
def slides_to_images(ppt_path, output_folder):
    subprocess.run([
        '/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'pdf', ppt_path, '--outdir', output_folder
    ], check=True)
    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    return [img.save(os.path.join(output_folder, f"slide_{i}.png"), 'PNG') or os.path.join(output_folder, f"slide_{i}.png")
            for i, img in enumerate(convert_from_path(pdf_path, dpi=200))]

# Step 6: Generate Video

def create_video(slide_imgs, audio_path, output_path):
    audio = AudioFileClip(audio_path)
    duration = audio.duration / len(slide_imgs)
    clips = [ImageClip(p).set_duration(duration) for p in slide_imgs]
    video = concatenate_videoclips(clips, method="compose").set_audio(audio)
    video.write_videofile(output_path, fps=24)

# Main

class Args:
    pdf = '/content/LLM_Overview.pdf'  # <-- Update with your actual PDF path
    avatar = '/content/man.png'  # <-- Update with your avatar image
    music = '/contents/breath-of-life_10-minutes-320859.mp3'  # Optional
    theme = 'creative'
    language = 'en'
    voice = 'enthusiastic'
    output = '/content/output/'



def main(pdf_path):

    args = Args()
    # Create configuration
    config = VideoConfig(
        theme=args.theme,
        language=args.language,
        voice_style=args.voice,
        include_background_music=bool(args.music)
    )
    text = extract_text_from_pdf(pdf_path)
    print("✅ Extracted text")

    chunks = chunk_text(text)
    print(f"✅ Split into {len(chunks)} chunks")

    results = [generate_chunk_content(chunk,config) for chunk in chunks]
    print("✅ Generated slide + script chunks")

    full_script = "\n\n".join(r.voice_over_script for r in results)
    audio_file = "voice.mp3"
    generate_audio(full_script, audio_file)
    print("✅ Audio saved")

    ppt_file = "presentation.pptx"
    generate_presentation(results, ppt_file, config)
    print("✅ Slides created")

    with tempfile.TemporaryDirectory() as tmpdir:
        imgs = slides_to_images(ppt_file, tmpdir)
        create_video(imgs, audio_file, "final_video.mp4")
        print("✅ Video exported")

if __name__ == "__main__":
    main("./contents/Redux_in_React.pdf")
