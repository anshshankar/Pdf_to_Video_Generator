import fitz  # PyMuPDF
from openai import OpenAI
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from moviepy.editor import AudioFileClip, ImageClip, concatenate_videoclips
from pdf2image import convert_from_path
from pydantic import BaseModel, ValidationError
import subprocess
import tempfile
import os
import json
import textwrap
from dotenv import load_dotenv

load_dotenv()

client = OpenAI(
    base_url=os.getenv("ENDPOINT"),
    api_key=os.getenv("TOKEN"),
)

# Pydantic schema for validation
class SlideChunk(BaseModel):
    slides: list[str]
    voice_over_script: str

# Step 1: Extract Text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return "".join(page.get_text() for page in doc)

# Step 2: Chunk and summarize text for slides and voiceover
def chunk_text(text, max_chunk_chars=2500):
    paragraphs = textwrap.wrap(text, max_chunk_chars, break_long_words=False, replace_whitespace=False)
    return paragraphs

def generate_chunk_content(chunk):
    prompt = (
        "Return JSON with two fields: 'slides' (list of concise bullet points summarizing the content) "
        "and 'voice_over_script' (a detailed voice-over narration). "
        "Format: {\"slides\": [...], \"voice_over_script\": \"...\"}. "
        "Only output raw JSON, no markdown formatting like triple backticks.\n\n"
        f"Content:\n{chunk}\n\nJSON Output:"
    )
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
        max_tokens=4000
    ).choices[0].message.content.strip()

    if response.startswith("```json"):
        response = response.lstrip("```json").rstrip("```").strip()
    elif response.startswith("```"):
        response = response.lstrip("```").rstrip("```").strip()

    try:
        parsed = json.loads(response)
        validated = SlideChunk(**parsed)
    except (json.JSONDecodeError, ValidationError) as e:
        raise ValueError(f"GPT format error: {e}\nRaw: {response}")

    return validated

# Step 3: Generate Audio from Script
def generate_audio(script, output_file):
    tts = gTTS(script)
    tts.save(output_file)

# Step 4: Generate Multi-Slide Presentation
def generate_presentation(slide_chunks, ppt_file):
    prs = Presentation()
    for i, chunk in enumerate(slide_chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i+1}"
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for point in chunk.slides:
            p = frame.add_paragraph()
            p.text = point
            p.font.size = Pt(24)
            p.space_after = Pt(12)
    prs.save(ppt_file)

# Step 5: Convert Slides to Images
def slides_to_images(ppt_path, output_folder):
    subprocess.run([
        '/Applications/LibreOffice.app/Contents/MacOS/soffice', '--headless', '--convert-to', 'pdf', ppt_path, '--outdir', output_folder
    ], check=True)
    pdf_path = os.path.join(output_folder, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    return [img.save(os.path.join(output_folder, f"slide_{i}.png"), 'PNG') or os.path.join(output_folder, f"slide_{i}.png")
            for i, img in enumerate(convert_from_path(pdf_path, dpi=200))]

# Step 6: Generate Video

def create_video(slide_imgs, audio_path, output_path):
    audio = AudioFileClip(audio_path)
    duration = audio.duration / len(slide_imgs)
    clips = [ImageClip(p).set_duration(duration) for p in slide_imgs]
    video = concatenate_videoclips(clips, method="compose").set_audio(audio)
    video.write_videofile(output_path, fps=24)

# Main

def main(pdf_path):
    text = extract_text_from_pdf(pdf_path)
    print("✅ Extracted text")

    chunks = chunk_text(text)
    print(f"✅ Split into {len(chunks)} chunks")

    results = [generate_chunk_content(chunk) for chunk in chunks]
    print("✅ Generated slide + script chunks")

    full_script = "\n\n".join(r.voice_over_script for r in results)
    audio_file = "voice.mp3"
    generate_audio(full_script, audio_file)
    print("✅ Audio saved")

    ppt_file = "presentation.pptx"
    generate_presentation(results, ppt_file)
    print("✅ Slides created")

    with tempfile.TemporaryDirectory() as tmpdir:
        imgs = slides_to_images(ppt_file, tmpdir)
        create_video(imgs, audio_file, "final_video.mp4")
        print("✅ Video exported")

if __name__ == "__main__":
    main("./contents/Redux_in_React.pdf")
